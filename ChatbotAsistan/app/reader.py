from __future__ import annotations
from pathlib import Path
from typing import Iterable, List, Dict, Callable, Tuple
import re

def _clean(s: str) -> str:
    s = s.replace("\x00", " ").replace("\ufeff", " ")
    s = re.sub(r"\s+", " ", s).strip()
    return s

def _yield_nonempty(lines: Iterable[str]) -> Iterable[str]:
    for ln in lines:
        t = _clean(ln)
        if t:
            yield t


def read_docx(path: Path) -> List[Dict]:
    from docx import Document
    doc = Document(path)
    out = []
    cid = 0
    for p in doc.paragraphs:
        t = _clean(p.text)
        if not t: continue
        out.append({"source": path.name, "page": None, "chunk_id": cid, "text": t})
        cid += 1
    return out


def read_pdf(path: Path) -> List[Dict]:
    from pypdf import PdfReader
    out = []
    cid = 0
    pdf = PdfReader(path.open("rb"))
    for i, page in enumerate(pdf.pages, start=1):
        text = page.extract_text() or ""
        for line in _yield_nonempty(text.splitlines()):
            out.append({"source": path.name, "page": i, "chunk_id": cid, "text": line})
            cid += 1
    return out

def read_pptx(path: Path) -> List[Dict]:
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    cid = 0
    for i, slide in enumerate(prs.slides, start=1):
        # başlık + tüm şekillerdeki metin kutuları
        texts = []
        if slide.shapes.title and slide.shapes.title.text:
            texts.append(slide.shapes.title.text)
        for sh in slide.shapes:
            if hasattr(sh, "text") and sh.text:
                texts.append(sh.text)
        block = " • ".join(_clean(t) for t in texts if _clean(t))
        for line in _yield_nonempty(block.split("•")):
            out.append({"source": path.name, "page": i, "chunk_id": cid, "text": line})
            cid += 1
    return out


def read_html(path: Path) -> List[Dict]:
    from bs4 import BeautifulSoup
    html = path.read_text("utf-8", errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    # görünür metni al: başlıklar + paragraflar + li
    parts = []
    for tag in soup.select("h1,h2,h3,h4,h5,h6,p,li"):
        t = _clean(tag.get_text(" ", strip=True))
        if t:
            parts.append(t)
    out = []
    cid = 0
    for line in parts:
        out.append({"source": path.name, "page": None, "chunk_id": cid, "text": line})
        cid += 1
    return out


def read_txt(path: Path) -> List[Dict]:
    out, cid = [], 0
    with path.open("r", encoding="utf-8", errors="ignore") as f:
        for line in _yield_nonempty(f):
            out.append({"source": path.name, "page": None, "chunk_id": cid, "text": line})
            cid += 1
    return out

ReaderFn = Callable[[Path], List[Dict]]
READERS: Dict[str, ReaderFn] = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".html": read_html,
    ".htm": read_html,
    ".txt": read_txt,
}

def load_dir(hrdir, allow_ext: Tuple[str, ...] = ()) -> List[Dict]:
    hrdir = Path(hrdir)
    hrdir.mkdir(parents=True, exist_ok=True)

    # allow_ext boşsa READERS.keys() kullan (mevcut mantık)
    exts = set(allow_ext) if allow_ext else set(READERS.keys())

    items: List[Dict] = []
    for p in hrdir.glob("*"):
        if p.is_file() and p.suffix.lower() in exts:
            try:
                out = READERS[p.suffix.lower()](p)
                # Reader her zaman liste döndürmeyebilir; normalize et
                if isinstance(out, dict):
                    items.append(out)
                elif isinstance(out, list):
                    items.extend(out)
                else:
                    print(f"[WARN] Reader {p.name} beklenmeyen tip döndürdü: {type(out)}")
            except Exception as e:
                print(f"[WARN] Okunamadı: {p.name} ({e})")
    return items
